from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
import os


# Usage 
# # Add Title Slide
# builder.add_title_slide(title="Model Explainability Report", subtitle="Customer Churn Prediction")
# # Add Global Summary Plot Slide
# builder.add_plot_slide(title="SHAP Summary Plot", plot_path=summary_plot_path, narrative="This plot shows overall feature impacts...")

class SlideDeckBuilder:
    def __init__(self, output_file: str = "Explainability_Report.pptx"):
        self.prs = Presentation()
        self.output_file = output_file

    def add_title_slide(self, title: str, subtitle: str = ""):
        slide_layout = self.prs.slide_layouts[0]  # Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_plot_slide(self, title: str, plot_path: str, narrative: str = ""):
        slide_layout = self.prs.slide_layouts[5]  # Title and Content Layout
        slide = self.prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title

        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(3.5)

        # Insert plot image
        if os.path.exists(plot_path):
            slide.shapes.add_picture(plot_path, left, top, height=height)

        # Add narrative below or beside
        if narrative:
            left_text = Inches(5.5)
            top_text = Inches(1.5)
            width_text = Inches(4)
            height_text = Inches(3.5)
            textbox = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = narrative
            p.font.size = Pt(14)

    def save_presentation(self):
        self.prs.save(self.output_file)
        print(f"✅ Presentation saved at {self.output_file}")
